#!/usr/bin/env python3
"""
STARK Procurement Scope Brief - PowerPoint Generator
Creates an on-brand presentation for Kim Christensen
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# STARK Brand Colors
STARK_NAVY = RGBColor(0x00, 0x32, 0x6E)
STARK_ORANGE = RGBColor(0xF5, 0x82, 0x1E)
STARK_LIGHT_LAVENDER = RGBColor(0xB9, 0xBD, 0xD7)
STARK_BLUE = RGBColor(0x0A, 0x5A, 0xC8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Paths
PNG_DIR = "/Users/dhk/Projects/STARK_Procurement/app/docs/architecture/drawio/png"
OUTPUT_PATH = "/Users/dhk/Projects/STARK_Procurement/app/docs/STARK_Procurement_Scope_Brief_Kim_Christensen.pptx"

def add_title_slide(prs, title, subtitle):
    """Add a title slide with STARK branding"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Navy background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = STARK_NAVY
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.5), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.5), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = STARK_LIGHT_LAVENDER
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    # Orange accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4), Inches(5.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = STARK_ORANGE
    line.line.fill.background()

    return slide

def add_section_slide(prs, title):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Navy background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = STARK_NAVY
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.5), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, image_path=None):
    """Add a content slide with optional image"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # White background (default)

    # Navy header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = STARK_NAVY
    header.line.fill.background()

    # Title in header
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"

    # Content area
    if image_path and os.path.exists(image_path):
        # Image on right, bullets on left
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
        slide.shapes.add_picture(image_path, Inches(6.2), Inches(1.4), width=Inches(7))
    else:
        # Full width bullets
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.5), Inches(5.5))

    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = STARK_NAVY
        p.font.name = "Arial"
        p.space_before = Pt(8)
        p.level = 0

    return slide

def add_diagram_slide(prs, title, image_path, caption=""):
    """Add a full-width diagram slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Navy header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = STARK_NAVY
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"

    # Diagram image - maximized
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.3), Inches(1.1), width=Inches(12.9))

    # Caption
    if caption:
        caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12.5), Inches(0.4))
        tf = caption_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = STARK_NAVY
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a data table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Navy header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = STARK_NAVY
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"

    # Table
    cols = len(headers)
    table_rows = len(rows) + 1
    table = slide.shapes.add_table(table_rows, cols, Inches(0.5), Inches(1.3), Inches(12.5), Inches(5.5)).table

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = STARK_NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.size = Pt(14)
        p.font.name = "Arial"

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF7)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = STARK_NAVY
            p.font.name = "Arial"

    return slide

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # === TITLE SLIDE ===
    add_title_slide(
        prs,
        "STARK Procurement",
        "Scope Definition Brief for Kim Christensen  •  March 2026"
    )

    # === EXECUTIVE SUMMARY ===
    add_section_slide(prs, "Executive Summary")

    add_content_slide(prs, "What We Are Building", [
        "A Procurement System handling the complete Purchase-to-Pay workflow:",
        "",
        "  PR Ingestion → PO Generation → Supplier Communication → Invoice Matching → Approval",
        "",
        "NOT building: Warehouse Management, Financial Accounting, E-commerce, Demand Planning"
    ])

    add_table_slide(prs, "Scale & Volume",
        ["Metric", "Annual Volume"],
        [
            ["Purchase Orders", "750,000"],
            ["PO Value", "€1.1B"],
            ["Suppliers", "10,000"],
            ["Branches", "84"],
            ["System Users", "30 procurement staff"],
        ]
    )

    # === SYSTEM CONTEXT ===
    add_section_slide(prs, "System Architecture")

    add_diagram_slide(
        prs,
        "System Context — External Boundaries",
        f"{PNG_DIR}/01-system-context.png",
        "C4 Level 1: STARK Procurement and its interactions with external systems"
    )

    add_diagram_slide(
        prs,
        "Container Architecture — Internal Components",
        f"{PNG_DIR}/02-container-diagram.png",
        "C4 Level 2: Web Application, API Layer, Database, Cache, Kafka Consumers/Producers"
    )

    # === PROCESS & WORKFLOW ===
    add_section_slide(prs, "Process & Workflow")

    add_diagram_slide(
        prs,
        "End-to-End Process Flow",
        f"{PNG_DIR}/03-process-flow.png",
        "PR to Payment: 6 stages from ingestion through invoice approval"
    )

    add_diagram_slide(
        prs,
        "Approval Workflow",
        f"{PNG_DIR}/05-approval-workflow.png",
        "Authority matrix, delegation, multi-level approval (4-eyes), and escalation"
    )

    # === INTEGRATION ===
    add_section_slide(prs, "Integration Architecture")

    add_diagram_slide(
        prs,
        "Kafka Integration — Topic Topology",
        f"{PNG_DIR}/04-kafka-integration.png",
        "Per-source inbound topics, per-channel outbound, DLQ patterns"
    )

    add_content_slide(prs, "Kafka-Native Integration", [
        "• Per-source inbound topics: Relex, ECom, SalesApp, Aspect4",
        "• Per-channel outbound topics: EDI, Email, Finance exports",
        "• Dead Letter Queues (DLQ) for failure handling per source/channel",
        "• Event topics for informational notifications (IDs + metadata)",
        "• REST fallback for non-Kafka sources (internally publishes to Kafka)",
        "• ACL-based security with SASL/mTLS authentication",
        "• Customer provisions topics; vendor delivers specs and code"
    ])

    # === DATA MODEL ===
    add_section_slide(prs, "Data Model")

    add_diagram_slide(
        prs,
        "Entity Relationships",
        f"{PNG_DIR}/06-data-model.png",
        "Core entities: PR, PO, Invoice, Approval, Delegation, Supplier"
    )

    # === USER EXPERIENCE ===
    add_section_slide(prs, "User Experience")

    add_table_slide(prs, "Users & Roles",
        ["Role", "Count", "Primary Function"],
        [
            ["Buyer", "~20", "PR review, PO creation, bundling, invoice matching"],
            ["Approver", "~8", "PO approval, invoice approval, exception handling"],
            ["Admin", "~2", "System configuration, user support, reporting"],
        ]
    )

    add_diagram_slide(
        prs,
        "Screen Navigation — 16 Screens",
        f"{PNG_DIR}/07-screen-navigation.png",
        "Dashboard-centric navigation with functional modules"
    )

    add_diagram_slide(
        prs,
        "Buyer's Morning — User Journey",
        f"{PNG_DIR}/08-user-journey.png",
        "Pain point identified: 'Resolve discrepancies' needs careful UX design"
    )

    # === SCOPE SUMMARY ===
    add_section_slide(prs, "Scope Summary")

    add_content_slide(prs, "What's IN Scope", [
        "• 16 screens: Dashboard, PR Inbox, PO Management, Invoice Matching, Approvals",
        "• Bundling Workspace: Group PRs into optimized POs with packet labeling",
        "• 2-Way Invoice Match: PO to Invoice with tolerance rules",
        "• Approval Workflow: Per-user limits, multi-level chains, delegation",
        "• Kafka Integration: Per-source consumers, per-channel producers, DLQ handling",
        "• Audit Trail: Full traceability of all decisions and changes"
    ])

    add_content_slide(prs, "What's OUT of Scope (Phase 1)", [
        "• Stock Transfer Orders (STO) — Requires NYCE WMS [FUTURE]",
        "• 3-Way Invoice Match — Requires goods receipt from NYCE [FUTURE]",
        "• SAP Finance Integration — Direct API [FUTURE]",
        "• Supplier Portal — Enhancement, not MVP [Phase 3]",
        "• Mobile Native App — Enhancement [Phase 3]",
        "• Advanced Analytics/BI — Different project [BI team]"
    ])

    add_content_slide(prs, "Dependencies We Assume Ready", [
        "• Relex PR feed — READY (automated replenishment)",
        "• ECom PR feed — READY (drop-shipment PRs)",
        "• Stark Output — READY (EDI/Email delivery to suppliers)",
        "• SSO/IAM — READY (user authentication)",
        "• STARK Kafka cluster — READY (customer provisions topics)",
        "• Supplier Master Data — READY (via MDM)"
    ])

    # === CLOSING ===
    add_section_slide(prs, "Next Steps")

    add_content_slide(prs, "Proposed Next Steps", [
        "1. Review and sign-off on scope boundaries",
        "2. Confirm Kafka topic provisioning timeline",
        "3. Align on integration schema specifications",
        "4. Schedule technical deep-dive with STARK IT team",
        "5. Finalize fixed-price commitment based on agreed scope"
    ])

    # Final slide
    add_title_slide(
        prs,
        "Thank You",
        "Questions & Discussion"
    )

    # Save
    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to: {OUTPUT_PATH}")
    return OUTPUT_PATH

if __name__ == "__main__":
    create_presentation()
