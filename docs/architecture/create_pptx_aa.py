#!/usr/bin/env python3
"""
STARK Procurement Scope Brief - Agentic Agency On-Brand Version
Black + Cement (#E6E6E1) palette, brutalist hard shadows aesthetic
For Kim Christensen
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Agentic Agency Brand Colors
AA_BLACK = RGBColor(0x00, 0x00, 0x00)
AA_CEMENT = RGBColor(0xE6, 0xE6, 0xE1)
AA_DARK_BG = RGBColor(0x11, 0x11, 0x11)

# Opacity simulation (pre-mixed with black background)
AA_CEMENT_90 = RGBColor(0xCF, 0xCF, 0xCA)
AA_CEMENT_70 = RGBColor(0xA1, 0xA1, 0x9D)
AA_CEMENT_50 = RGBColor(0x73, 0x73, 0x70)
AA_CEMENT_30 = RGBColor(0x45, 0x45, 0x43)
AA_CEMENT_20 = RGBColor(0x2E, 0x2E, 0x2C)

# Paths
PNG_DIR = "/Users/dhk/Projects/STARK_Procurement/app/docs/architecture/drawio/png"
OUTPUT_PATH = "/Users/dhk/Projects/STARK_Procurement/app/docs/STARK_Procurement_Scope_Brief_Kim_Christensen.pptx"

def add_hard_shadow(slide, left, top, width, height, offset=8):
    """Add a brutalist hard shadow (offset rectangle)"""
    shadow = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left + offset/72),
        Inches(top + offset/72),
        Inches(width),
        Inches(height)
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = AA_CEMENT_20
    shadow.line.fill.background()
    return shadow

def add_title_slide(prs, title, subtitle):
    """Add a title slide with AA branding - black bg, cement text"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Black background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = AA_BLACK
    background.line.fill.background()

    # Decorative geometric element (top right)
    geo1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10), Inches(0.5), Inches(2), Inches(2))
    geo1.fill.background()
    geo1.line.color.rgb = AA_CEMENT_30
    geo1.line.width = Pt(1)
    geo1.rotation = 12

    # Decorative line (bottom left)
    line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6), Inches(3), Inches(0.02))
    line1.fill.solid()
    line1.fill.fore_color.rgb = AA_CEMENT_30
    line1.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = AA_CEMENT
    p.font.name = "Arial"  # Space Grotesk fallback
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.5), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = AA_CEMENT_70
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.LEFT

    # "Engineering > Prompting" tagline (bottom right)
    tagline = slide.shapes.add_textbox(Inches(8.5), Inches(6.8), Inches(4.5), Inches(0.5))
    tf = tagline.text_frame
    p = tf.paragraphs[0]
    p.text = "Engineering > Prompting"
    p.font.size = Pt(11)
    p.font.color.rgb = AA_CEMENT_50
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.RIGHT

    return slide

def add_section_slide(prs, title, number=None):
    """Add a section divider slide - minimal, bold"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Black background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = AA_BLACK
    background.line.fill.background()

    # Section number (if provided)
    if number:
        num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(2), Inches(1))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"0{number}"
        p.font.size = Pt(48)
        p.font.color.rgb = AA_CEMENT_30
        p.font.name = "Arial"

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = AA_CEMENT
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.LEFT

    # Horizontal line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.8), Inches(4), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = AA_CEMENT_50
    line.line.fill.background()

    return slide

def add_content_slide(prs, title, bullets, image_path=None):
    """Add a content slide with bullets"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Black background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = AA_BLACK
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = AA_CEMENT
    p.font.name = "Arial"

    # Underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(3), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = AA_CEMENT_30
    line.line.fill.background()

    # Content area
    if image_path and os.path.exists(image_path):
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(5))
        # Add image with hard shadow effect
        shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(1.85), Inches(6), Inches(4.5))
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = AA_CEMENT_20
        shadow.line.fill.background()
        slide.shapes.add_picture(image_path, Inches(6.5), Inches(1.7), width=Inches(6))
    else:
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.5))

    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle empty lines
        if bullet.strip() == "":
            p.text = ""
            p.space_before = Pt(6)
            continue

        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = AA_CEMENT_90 if not bullet.startswith("  ") else AA_CEMENT_70
        p.font.name = "Arial"
        p.space_before = Pt(12)
        p.level = 0

    return slide

def add_diagram_slide(prs, title, image_path, caption=""):
    """Add a full-width diagram slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Black background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = AA_BLACK
    background.line.fill.background()

    # Title (small, top left)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = AA_CEMENT
    p.font.name = "Arial"

    # Diagram container with hard shadow
    if os.path.exists(image_path):
        # Shadow
        shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.15), Inches(12.2), Inches(5.6))
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = AA_CEMENT_20
        shadow.line.fill.background()

        # White container for diagram
        container = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(12.2), Inches(5.6))
        container.fill.solid()
        container.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        container.line.fill.background()

        # Image
        slide.shapes.add_picture(image_path, Inches(0.6), Inches(1.1), width=Inches(12))

    # Caption
    if caption:
        caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.5), Inches(0.5))
        tf = caption_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = AA_CEMENT_50
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.LEFT

    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a data table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Black background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = AA_BLACK
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = AA_CEMENT
    p.font.name = "Arial"

    # Table shadow
    cols = len(headers)
    table_rows = len(rows) + 1
    shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(1.75), Inches(11.5), Inches(0.6 * table_rows))
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = AA_CEMENT_20
    shadow.line.fill.background()

    # Table
    table = slide.shapes.add_table(table_rows, cols, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.6 * table_rows)).table

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = AA_CEMENT
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = AA_BLACK
        p.font.size = Pt(14)
        p.font.name = "Arial"

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = AA_DARK_BG if row_idx % 2 == 0 else AA_BLACK
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = AA_CEMENT_90
            p.font.name = "Arial"

    return slide

def add_quote_slide(prs, quote, attribution=""):
    """Add a quote slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Black background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = AA_BLACK
    background.line.fill.background()

    # Left border line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(0.08), Inches(3))
    line.fill.solid()
    line.fill.fore_color.rgb = AA_CEMENT_50
    line.line.fill.background()

    # Quote
    quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10.5), Inches(2.5))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'"{quote}"'
    p.font.size = Pt(32)
    p.font.italic = True
    p.font.color.rgb = AA_CEMENT_90
    p.font.name = "Arial"

    # Attribution
    if attribution:
        attr_box = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(10.5), Inches(0.5))
        tf = attr_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"— {attribution}"
        p.font.size = Pt(14)
        p.font.color.rgb = AA_CEMENT_50
        p.font.name = "Arial"

    return slide

def create_presentation():
    """Create the full AA-branded presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # === TITLE SLIDE ===
    add_title_slide(
        prs,
        "STARK Procurement",
        "Scope Definition Brief  •  Kim Christensen  •  March 2026"
    )

    # === EXECUTIVE SUMMARY ===
    add_section_slide(prs, "Executive Summary", 1)

    add_content_slide(prs, "What We're Building", [
        "A Procurement System handling the complete Purchase-to-Pay workflow:",
        "",
        "PR Ingestion → PO Generation → Supplier Communication → Invoice Matching → Approval",
        "",
        "NOT building:",
        "  Warehouse Management • Financial Accounting • E-commerce • Demand Planning"
    ])

    add_table_slide(prs, "Scale & Volume",
        ["Metric", "Annual Volume"],
        [
            ["Purchase Orders", "750,000"],
            ["PO Value", "€1.1B"],
            ["Suppliers", "10,000"],
            ["Branches", "84"],
            ["System Users", "30 procurement staff"],
        ]
    )

    # === SYSTEM CONTEXT ===
    add_section_slide(prs, "System Architecture", 2)

    add_diagram_slide(
        prs,
        "System Context — External Boundaries",
        f"{PNG_DIR}/01-system-context.png",
        "C4 Level 1: STARK Procurement and its interactions with external systems"
    )

    add_diagram_slide(
        prs,
        "Container Architecture — Internal Components",
        f"{PNG_DIR}/02-container-diagram.png",
        "C4 Level 2: Web Application, API Layer, Database, Cache, Kafka Consumers/Producers"
    )

    # === PROCESS & WORKFLOW ===
    add_section_slide(prs, "Process & Workflow", 3)

    add_diagram_slide(
        prs,
        "End-to-End Process Flow",
        f"{PNG_DIR}/03-process-flow.png",
        "PR to Payment: 6 stages from ingestion through invoice approval"
    )

    add_diagram_slide(
        prs,
        "Approval Workflow",
        f"{PNG_DIR}/05-approval-workflow.png",
        "Authority matrix, delegation, multi-level approval (4-eyes), and escalation"
    )

    # === INTEGRATION ===
    add_section_slide(prs, "Integration Architecture", 4)

    add_diagram_slide(
        prs,
        "Kafka Integration — Topic Topology",
        f"{PNG_DIR}/04-kafka-integration.png",
        "Per-source inbound topics, per-channel outbound, DLQ patterns"
    )

    add_content_slide(prs, "Kafka-Native Integration", [
        "Per-source inbound topics: Relex, ECom, SalesApp, Aspect4",
        "Per-channel outbound topics: EDI, Email, Finance exports",
        "Dead Letter Queues (DLQ) for failure handling per source/channel",
        "Event topics for informational notifications (IDs + metadata)",
        "REST fallback for non-Kafka sources",
        "ACL-based security with SASL/mTLS authentication",
        "Customer provisions topics; vendor delivers specs and code"
    ])

    # === DATA MODEL ===
    add_section_slide(prs, "Data Model", 5)

    add_diagram_slide(
        prs,
        "Entity Relationships",
        f"{PNG_DIR}/06-data-model.png",
        "Core entities: PR, PO, Invoice, Approval, Delegation, Supplier"
    )

    # === USER EXPERIENCE ===
    add_section_slide(prs, "User Experience", 6)

    add_table_slide(prs, "Users & Roles",
        ["Role", "Count", "Primary Function"],
        [
            ["Buyer", "~20", "PR review, PO creation, bundling, invoice matching"],
            ["Approver", "~8", "PO approval, invoice approval, exception handling"],
            ["Admin", "~2", "System configuration, user support, reporting"],
        ]
    )

    add_diagram_slide(
        prs,
        "Screen Navigation — 16 Screens",
        f"{PNG_DIR}/07-screen-navigation.png",
        "Dashboard-centric navigation with functional modules"
    )

    add_diagram_slide(
        prs,
        "Buyer's Morning — User Journey",
        f"{PNG_DIR}/08-user-journey.png",
        "Pain point identified: 'Resolve discrepancies' needs careful UX design"
    )

    # === SCOPE SUMMARY ===
    add_section_slide(prs, "Scope Summary", 7)

    add_content_slide(prs, "What's IN Scope", [
        "16 screens: Dashboard, PR Inbox, PO Management, Invoice Matching, Approvals",
        "Bundling Workspace: Group PRs into optimized POs with packet labeling",
        "2-Way Invoice Match: PO to Invoice with tolerance rules",
        "Approval Workflow: Per-user limits, multi-level chains, delegation",
        "Kafka Integration: Per-source consumers, per-channel producers, DLQ handling",
        "Audit Trail: Full traceability of all decisions and changes"
    ])

    add_content_slide(prs, "What's OUT of Scope (Phase 1)", [
        "Stock Transfer Orders (STO) — Requires NYCE WMS [FUTURE]",
        "3-Way Invoice Match — Requires goods receipt from NYCE [FUTURE]",
        "SAP Finance Integration — Direct API [FUTURE]",
        "Supplier Portal — Enhancement, not MVP [Phase 3]",
        "Mobile Native App — Enhancement [Phase 3]",
        "Advanced Analytics/BI — Different project [BI team]"
    ])

    add_content_slide(prs, "Dependencies We Assume Ready", [
        "Relex PR feed — READY (automated replenishment)",
        "ECom PR feed — READY (drop-shipment PRs)",
        "Stark Output — READY (EDI/Email delivery to suppliers)",
        "SSO/IAM — READY (user authentication)",
        "STARK Kafka cluster — READY (customer provisions topics)",
        "Supplier Master Data — READY (via MDM)"
    ])

    # === CLOSING ===
    add_section_slide(prs, "Next Steps", 8)

    add_content_slide(prs, "Proposed Next Steps", [
        "1. Review and sign-off on scope boundaries",
        "2. Confirm Kafka topic provisioning timeline",
        "3. Align on integration schema specifications",
        "4. Schedule technical deep-dive with STARK IT team",
        "5. Finalize fixed-price commitment based on agreed scope"
    ])

    # Quote slide
    add_quote_slide(
        prs,
        "There's a gap between 'AI helped me write this function' and 'AI systematically delivered this feature.' That gap has a name.",
        "AGENTIC ENGINEERING"
    )

    # Final slide
    add_title_slide(
        prs,
        "Thank You",
        "Questions & Discussion"
    )

    # Save
    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to: {OUTPUT_PATH}")
    return OUTPUT_PATH

if __name__ == "__main__":
    create_presentation()
